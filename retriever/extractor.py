import PyPDF2
import easyocr
import numpy as np
import re
from PIL import Image
from bs4 import BeautifulSoup
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.options import Options
import time
from docx import Document
from pptx import Presentation


class TextExtractor:
    def fromTextFile(self, file) -> str:
        file_content = file.read().decode("utf-8")
        return file_content

    def fromPDFFile(self, file) -> str:
        text = ""
        reader = PyPDF2.PdfReader(file)
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text += page.extract_text()
        return text

    def fromImageFileEasyOCR(self, image) -> str:
        text = ""
        pil_image = Image.open(image)
        image_array = np.asarray(pil_image)
        reader = easyocr.Reader(["en"], verbose=False)
        result = reader.readtext(image_array)
        for res in result:
            text += res[-2] + " "
        return text


    def from_ppt_file(self, file_path) -> str:
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
            return text.strip()
        except Exception as e:
            return f"Error reading PPT file: {e}"
        
    # def fromDocxFile(self, file) -> str:
    #     doc = Document(file)
    #     text = ""
    #     for para in doc.paragraphs:
    #         text += para.text
    #     return text

    def fromUrl(self, url: str) -> str | None:
        options = Options()
        options.add_argument("--headless")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-gpu")
        options.add_argument(
            "user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.36"
        )

        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=options)

        try:
            driver.get(url)
            time.sleep(5)

            soup = BeautifulSoup(driver.page_source, "html.parser")
            text = soup.get_text(separator=" ", strip=True)
            return text
        finally:
            driver.quit()
            
    def fromDocxFile(self, file) -> str:
        doc = Document(file)
        full_text = []
        for para in doc.paragraphs:
            html_paragraph = ""
            for run in para.runs:
                text = run.text
                if run.bold:
                    text = f"<b>{text}</b>"
                if run.italic:
                    text = f"<i>{text}</i>"
                html_paragraph += text
            full_text.append(html_paragraph)
        return "<br>".join(full_text)
    
    def fromHTMLFile(self, file) -> str:
        soup = BeautifulSoup(file, "html.parser")
        text = soup.get_text()
        text = re.sub(r"\s+", " ", text)
        # remove \n and \t
        text = re.sub(r"\n", " ", text)
        text = text.strip()
        return text


if __name__ == "__main__":
    extractor = TextExtractor()
    extractor.fromImageFileEasyOCR("car.png")
    print(extractor.fromTextFile(open("./car.txt")))
    print(extractor.fromPDFFile(open("car.pdf", "rb")))
    print(extractor.fromImageFile("car.png"))
    print(extractor.fromImageFileEasyOCR(open("car.png", "rb")))
    print(extractor.fromUrl("https://en.wikipedia.org/wiki/Animal"))
    print(extractor.fromHTMLFile(open("content.html")))
